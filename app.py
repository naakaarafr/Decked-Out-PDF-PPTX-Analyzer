import streamlit as st
import os
import io
from PIL import Image
import base64
import warnings
import google.generativeai as genai
import tempfile
from pptx import Presentation

# Configure Streamlit page with custom CSS for better note presentation
st.set_page_config(
    page_title="Decked Out PDF/PPTX Analyzer",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Load environment variables from .env file if present
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# Initialize session state for notes type selection, storing generated notes, and chat history
if 'selected_notes_type' not in st.session_state:
    st.session_state.selected_notes_type = None
if 'notes_content' not in st.session_state:
    st.session_state.notes_content = None
if 'note_type' not in st.session_state:
    st.session_state.note_type = None
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

# Define show_debug setting
show_debug = True  # Set to True to see more debugging information

# Set up pdf2image with poppler
try:
    from pdf2image import convert_from_bytes, convert_from_path
    PDF2IMAGE_AVAILABLE = True
    #st.sidebar.success("PDF2Image is available!")
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    #st.sidebar.warning("PDF2Image not available. For better results, install pdf2image and poppler.")

# Initialize Google Generative AI client
api_key = os.getenv("GOOGLE_API_KEY")
if api_key:
    genai.configure(api_key=api_key)
else:
    st.sidebar.error("Google API Key not found. Please set the GOOGLE_API_KEY environment variable.")

# Set model to use (from environment variable or default)
model_name = os.getenv("MODEL", "gemini-2.0-flash")

# Custom CSS for better formatting of notes
st.markdown("""
<style>
    /* Main styling */
    .main-header {
        font-size: 2.8rem;
        font-weight: bold;
        background: linear-gradient(90deg, #CF9FFF, #512DA8);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.5rem;
        padding-top: 1rem;
        text-align: center;
    }
    
    .tagline {
        font-size: 1.1rem;
        color: #555;
        margin-bottom: 2rem;
    }
    
    .sub-header {
        font-size: 1.6rem;
        font-weight: 600;
        color: #1E88E5;
        margin-top: 1.5rem;
        margin-bottom: 1rem;
        padding-bottom: 0.3rem;
        border-bottom: 2px solid #f0f0f0;
    }
    
    /* Card styling */
    .card {
        background-color: white;
        border-radius: 10px;
        padding: 1.5rem;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        margin-bottom: 1.5rem;
        border: 1px solid #f0f0f0;
    }
    
    /* Notes styling */
    .notes-container {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        border-left: 5px solid #1E88E5;
        margin-bottom: 1.5rem;
        box-shadow: 0 2px 4px rgba(0,0,0,0.04);
        max-height: 70vh;
        overflow-y: auto;
    }
    
    /* Document preview */
    .pdf-preview {
        border: 1px solid #eee;
        border-radius: 8px;
        padding: 1rem;
        margin-bottom: 1rem;
        background-color: white;
        text-align: center;
    }
    
    /* Buttons styling */
    .stButton button {
        border-radius: 8px !important;
        font-weight: 500 !important;
        padding: 0.5rem 1rem !important;
        transition: all 0.2s ease !important;
    }
    
    .stButton button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 4px 8px rgba(0,0,0,0.1) !important;
    }
    
    .note-type-button {
        width: 100%;
        padding: 1rem 0.5rem !important;
        border-radius: 8px !important;
        font-weight: 500 !important;
    }
    
        /* Download button */
    .download-btn {
        display: inline-block;
        padding: 0.6rem 1.2rem;
        background: linear-gradient(90deg, #CF9FFF, #3949AB);
        color: white !important; /* Force white text with !important */
        text-decoration: none;
        border-radius: 8px;
        font-weight: 500;
        text-align: center;
        margin: 1rem 0;
        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        transition: all 0.2s ease;
    }

    .download-btn:hover {
        box-shadow: 0 4px 10px rgba(0,0,0,0.15);
        transform: translateY(-2px);
        color: white !important; /* Maintain white text on hover */
    }

    .download-btn * {
        /* Ensure all child elements are also white */
        color: white !important;
    }
    
    /* Chat styling */
    .chat-message {
        padding: 1rem;
        border-radius: 8px;
        margin-bottom: 1rem;
        max-width: 85%;
    }
    
    .user-message {
        background-color: #E3F2FD;
        margin-left: auto;
        border-bottom-right-radius: 0;
        border: 1px solid #BBDEFB;
    }
    
    .bot-message {
        background-color: #F5F5F5;
        margin-right: auto;
        border-bottom-left-radius: 0;
        border: 1px solid #E0E0E0;
    }
    
    /* Status indicators */
    .status-indicator {
        padding: 0.3rem 0.6rem;
        border-radius: 20px;
        font-size: 0.8rem;
        font-weight: 500;
        display: inline-block;
        margin-right: 0.5rem;
    }
    
    .status-success {
        background-color: #E8F5E9;
        color: #2E7D32;
        border: 1px solid #C8E6C9;
    }
    
    .status-pending {
        background-color: #FFF8E1;
        color: #F57F17;
        border: 1px solid #FFECB3;
    }
    
    .status-info {
        background-color: #E3F2FD;
        color: #1565C0;
        border: 1px solid #BBDEFB;
    }
    
    /* Tabs styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 2px;
    }
    
    .stTabs [data-baseweb="tab"] {
        padding: 0.75rem 1.5rem;
        border-radius: 8px 8px 0 0;
    }
    
    .stTabs [data-baseweb="tab-highlight"] {
        background-color: #1E88E5;
    }
    
    /* Hide Streamlit branding */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    
    /* Key highlight and terms */
    .key-highlight {
        background-color: #ffff99;
        padding: 2px 4px;
        border-radius: 3px;
    }
    
    .key-term {
        font-weight: bold;
        color: #0d47a1;
    }
    
    /* Footer styling */
    .footer {
        padding: 1.5rem;
        border-top: 1px solid #eee;
        margin-top: 2rem;
        color: #666;
    }
    
    .footer-header {
        font-size: 1.2rem;
        font-weight: 600;
        color: #333;
        margin-bottom: 1rem;
    }
    
    /* Code blocks */
    code {
        padding: 0.2rem 0.4rem;
        background-color: #f5f5f5;
        border-radius: 4px;
        font-size: 0.9rem;
    }
    
    /* Progress bar */
    .stProgress > div > div {
        background-color: #1E88E5;
    }
</style>
""", unsafe_allow_html=True)

st.markdown('<div class="main-header"><h1>📄 Decked Out PDF/PPTX Analyzer</h1></div>', unsafe_allow_html=True)
st.markdown(
    """
    Upload a PDF or PPTX document and get concise, effective notes highlighting key information.
    Select the type of notes you want based on your needs, then chat about them in the Chat tab.
    """
)

# Create tabs for PDF Notes and Chat
st.markdown('---', unsafe_allow_html=True)
tab1, tab2 = st.tabs(["PDF Notes", "Chat"])

# PDF Notes Tab - contains all original functionality with added PPTX support via text extraction
with tab1:
    st.markdown('<div class="sub-header">Upload Your Document</div>', unsafe_allow_html=True)
    file_types = ["pdf", "pptx"]
    uploaded_file = st.file_uploader("Upload your PDF or PPTX document:", type=file_types)

    # Define callback functions to update session state when buttons are clicked
    def set_official_notes():
        st.session_state.selected_notes_type = "official"

    def set_english_notes():
        st.session_state.selected_notes_type = "english"

    def set_hinglish_notes():
        st.session_state.selected_notes_type = "hinglish"

    # Function to extract text from PDF images using pdf2image (Poppler-based)
    def extract_pdf_text_with_poppler(pdf_bytes):
        """Extract text representation using pdf2image (Poppler-based) to convert PDF to images"""
        if not PDF2IMAGE_AVAILABLE:
            raise ImportError("pdf2image and poppler are required but not installed")
        
        try:
            # Convert PDF to images with additional parameters to handle problematic PDFs
            images = convert_from_bytes(
                pdf_bytes,
                dpi=200,  # Higher DPI to ensure better quality
                fmt='jpeg',  # Explicitly set format
                thread_count=1,  # Single-threaded for better stability
                strict=False,  # Less strict parsing
                use_cropbox=True,  # Use cropbox instead of mediabox
                transparent=False  # No transparency
            )
            
            if not images or len(images) == 0:
                raise Exception("No images extracted from PDF. The document may be empty or corrupted.")
            
            # For text extraction purposes, we'll describe the visual content
            extracted_content = ""
            
            for i, image in enumerate(images):
                # Save image to buffer
                img_byte_arr = io.BytesIO()
                image.save(img_byte_arr, format='JPEG')
                img_byte_arr = img_byte_arr.getvalue()
                
                # Add page marker
                extracted_content += f"\n\n--- Page {i + 1} ---\n\n"
                # We're not doing OCR here, just using the images directly
                extracted_content += f"[PDF Page {i + 1} converted to image]"
            
            return extracted_content, images
        except Exception as e:
            raise Exception(f"Error extracting with pdf2image: {str(e)}")

    # Function to extract text from PPTX files
    def extract_pptx_text(pptx_path):
        prs = Presentation(pptx_path)
        text_content = ""
        for i, slide in enumerate(prs.slides, 1):
            text_content += f"\n\n--- Slide {i} ---\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + "\n"
        return text_content

    # Function to process the document (PDF or PPTX)
    def input_document_setup(uploaded_file):
        """Process PDF or PPTX using appropriate extraction methods"""
        if uploaded_file is not None:
            try:
                file_ext = os.path.splitext(uploaded_file.name)[1].lower()
                
                if file_ext == '.pdf':
                    pdf_bytes = uploaded_file.read()
                    description, images = extract_pdf_text_with_poppler(pdf_bytes)
                    extraction_method = "pdf2image (Poppler)"
                    return description, extraction_method, images
                
                elif file_ext == '.pptx':
                    # Save PPTX to temporary file
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_pptx:
                        tmp_pptx.write(uploaded_file.read())
                        tmp_pptx_path = tmp_pptx.name
                    try:
                        text_content = extract_pptx_text(tmp_pptx_path)
                        extraction_method = "python-pptx text extraction"
                        return text_content, extraction_method, None
                    finally:
                        os.remove(tmp_pptx_path)
                
                else:
                    raise Exception("Unsupported file type. Please upload a PDF or PPTX file.")
            
            except Exception as e:
                st.error(f"Error processing document: {str(e)}")
                # Fallback to PyPDF2 for PDFs if applicable
                if file_ext == '.pdf':
                    try:
                        import PyPDF2
                        from io import BytesIO
                        
                        uploaded_file.seek(0)
                        pdf_reader = PyPDF2.PdfReader(BytesIO(uploaded_file.read()))
                        text_content = ""
                        
                        for i, page in enumerate(pdf_reader.pages):
                            page_text = page.extract_text()
                            if page_text:
                                text_content += f"\n\n--- Page {i + 1} ---\n\n"
                                text_content += page_text
                            else:
                                text_content += f"\n\n--- Page {i + 1} ---\n\n"
                                text_content += f"[PDF Page {i + 1} - No text extracted]"
                        
                        if not text_content.strip():
                            text_content = "No text could be extracted from the PDF."
                        
                        return text_content, "PyPDF2 text extraction", None
                    except ImportError:
                        return f"PDF extraction failed: {str(e)}. Install PyPDF2 for text extraction.", "Failed", None
                    except Exception as pdf_error:
                        return f"Extraction failed: {str(pdf_error)}", "Failed", None
                else:
                    return f"Error processing PPTX: {str(e)}", "Failed", None
        else:
            raise FileNotFoundError("No file uploaded")

    # Display document preview when uploaded (only for PDFs)
    if uploaded_file is not None:
        st.markdown('---', unsafe_allow_html=True)
        st.write("✅ Document uploaded successfully!")
        
        # Try to display first page preview for PDFs
        file_ext = os.path.splitext(uploaded_file.name)[1].lower()
        if file_ext == '.pdf' and PDF2IMAGE_AVAILABLE:
            try:
                pos = uploaded_file.tell()
                uploaded_file.seek(0)
                pdf_bytes = uploaded_file.read()
                images = convert_from_bytes(
                    pdf_bytes, 
                    first_page=1, 
                    last_page=1,
                    dpi=200,
                    fmt='jpeg',
                    thread_count=1,
                    strict=False,
                    use_cropbox=True,
                    transparent=False
                )
                if images:
                    st.image(images[0], width=300, caption="Preview of first page")
                uploaded_file.seek(pos)
            except Exception as e:
                st.info(f"Preview not available: {str(e)}")
        elif file_ext == '.pptx':
            st.info("Preview not available for PPTX files.")
        
        st.markdown('</div>', unsafe_allow_html=True)

    # Enhanced prompts for better key terms and important sentences extraction
    prompt_official_notes = """
    You are a professional document-analysis assistant specializing in key information extraction. Analyze the provided PDF and generate concise, structured official notes that highlight ALL the most important keywords, terms, and sentences. Your notes must:

    1. **Tone & Style**  
       • Formal, precise, and professional
       • Academic/business style with specialized terminology preserved exactly as written
       • Use complete sentences with proper grammar

    2. **Structure**  
       • Begin with a one-paragraph "Executive Summary" capturing the document's main purpose and key findings
       • Organize content under clear hierarchical headings that match the document's structure
       • Use numbered lists for processes/sequences, bullet points for discrete facts
       • Include page references in [brackets] for important information

    3. **Content Requirements**  
       • Extract and highlight ALL key terms, definitions, and technical vocabulary (in **bold**)
       • Include critical sentences verbatim where they contain essential information
       • Preserve important numerical data, statistics, and measurements exactly as presented
       • Identify and highlight main arguments, evidence, and conclusions
       • Note any actionable items, recommendations, or future directions

    4. **Emphasis**
       • **Bold** all key terms, crucial phrases, and technical vocabulary
       • Underline or quote particularly important complete sentences
       • Maintain the hierarchical relationship between concepts

    5. **Length & Completeness**  
       • 300-500 words total (comprehensive but concise)
       • Ensure ALL major concepts and terms are included, even if brief
       • Prioritize breadth of coverage over depth of explanation

    Produce the output as markdown text with appropriate formatting for headings, lists, and emphasis. Focus on creating a professional reference document that captures ALL essential information.
    """

    prompt_english_notes = """
    You are an expert document summarizer specializing in extracting key information in plain language. Analyze the provided PDF and generate clear, accessible English notes that capture ALL important concepts, terms, and sentences. Your notes must:

    1. **Tone & Style**  
       • Conversational and easy to understand for general audiences
       • Explain technical concepts in simple terms but preserve important terminology
       • Use short, direct sentences with clear meaning

    2. **Structure**  
       • Begin with a brief "Overview" (2-3 sentences) capturing the main purpose and takeaways
       • Organize under simple, descriptive headings that guide the reader
       • Use bullet points extensively to break down complex ideas
       • Number any steps, processes, or sequences

    3. **Content Requirements**  
       • Identify and **bold** ALL key terms, technical vocabulary, and specialized concepts
       • Extract important complete sentences that contain critical information (in "quotes")
       • Simplify complex ideas without losing essential meaning
       • Include any important numbers, measurements, or data points
       • Highlight practical applications or real-world implications

    4. **Emphasis**
       • **Bold** important terms and phrases throughout
       • Place particularly important sentences in "quotes"
       • Use simple formatting to visually separate different types of information

    5. **Length & Comprehensiveness**  
       • 250-400 words (concise but thorough)
       • Cover ALL major points from the document
       • Prioritize breadth to ensure nothing important is missed

    Produce the output as markdown text with appropriate formatting for headings, bullets, and emphasis. Focus on making the information accessible while ensuring ALL key words and important sentences are preserved.
    """

    prompt_hinglish_notes = """
    You are a friendly document summarizer specializing in creating accessible notes in mixed Hindi-English (Hinglish). Analyze the provided PDF and create notes that capture ALL important keywords, concepts and sentences in simple Hinglish using Roman script only. Your notes must:

    1. **Tone & Style**  
       • Conversational, simple, and friendly - jaise dost se baat kar rahe ho
       • Use a natural mix of Hindi and English - kuch technical terms English mein rakho
       • Short sentences and simple structure for easy understanding

    2. **Structure**  
       • "Overview" section mein 2-3 lines mein document ka main point batao
       • Simple headings jo content ko categorize karein
       • Har important point ko bullet points mein break karo
       • Steps ya process ko number karo

    3. **Content Requirements**  
       • Har **important keyword aur technical term ko bold** karo
       • Critical sentences ko "quotes" mein rakho, unki importance highlight karne ke liye
       • Difficult concepts ko everyday examples se explain karo
       • Important numbers, dates ya statistics ko exactly preserve karo
       • Document ke har major section se key points extract karo

    4. **Emphasis**
       • Important terms ko **bold** karo
       • Bahut important sentences ko "quotes" mein rakho
       • Different types of information ko visually separate karo

    5. **Length & Completeness**  
       • 250-400 words (concise par thorough)
       • Document ke SABHI major points cover karo
       • Har important keyword aur concept include karo

    Notes ko markdown text format mein banao with proper formatting for headings, bullets, and emphasis. Focus on making information accessible while ensuring ALL key terms and important sentences are preserved.
    """

    # Create columns for note type selection buttons 
    st.markdown('<div class="sub-header">Select Notes Type</div>', unsafe_allow_html=True)
# Define the radio button options
notes_options = ["Official Notes", "English Notes", "Hinglish Notes"]

# Define the descriptions for each option
notes_descriptions = {
    "Official Notes": "Formal, professional notes with technical terminology",
    "English Notes": "Simplified notes in plain English",
    "Hinglish Notes": "Notes in mixed Hindi-English for accessibility"
}

# Create the radio button with a horizontal layout
selected_note_type = st.radio(
    "Select notes type:",
    options=notes_options,
    horizontal=True,
    help="Choose your preferred notes format",
    key="notes_radio"
)

# Update session state based on selection
if selected_note_type == "Official Notes":
    st.session_state.selected_notes_type = "official"
elif selected_note_type == "English Notes":
    st.session_state.selected_notes_type = "english"
elif selected_note_type == "Hinglish Notes":
    st.session_state.selected_notes_type = "hinglish"

# Display current selection with description
if st.session_state.selected_notes_type:
    note_type_name = {
        "official": "Official Notes",
        "english": "English Notes",
        "hinglish": "Hinglish Notes"
    }.get(st.session_state.selected_notes_type)
    st.success(f"Selected: {note_type_name}")
    
    # Set selected prompt based on session state
    selected_prompt = None
    note_type = None

    if st.session_state.selected_notes_type == "official":
        selected_prompt = prompt_official_notes
        note_type = "Official Notes"
    elif st.session_state.selected_notes_type == "english":
        selected_prompt = prompt_english_notes
        note_type = "English Notes"
    elif st.session_state.selected_notes_type == "hinglish":
        selected_prompt = prompt_hinglish_notes
        note_type = "Hinglish Notes"

    # Function to create a download link for notes
    def get_download_link(notes_content, note_type):
        notes_bytes = notes_content.encode()
        b64 = base64.b64encode(notes_bytes).decode()
        filename = f"{note_type.replace(' ', '_').lower()}.md"
        href = f'<a href="data:text/markdown;base64,{b64}" download="{filename}" class="download-btn">Download {note_type}</a>'
        return href

    # Generate notes button
    if st.button("Generate Notes", key="generate_notes", help="Generate the selected type of notes"):
        if uploaded_file is not None:
            if not selected_prompt:
                st.warning("Please select what type of notes you want first.")
            else:
                with st.spinner(f"Analyzing document and generating {note_type}..."):
                    try:
                        # Process the document and prepare the input using extraction methods
                        description, extraction_method, images = input_document_setup(uploaded_file)
                        
                        # For debugging, show the extraction method
                        #if show_debug:
                            #st.info(f"Using {extraction_method} for extraction")
                        
                        # If we have images, display the first 3 for better analysis (only for PDFs)
                        if images and len(images) > 0:
                            st.subheader("Document Content Preview")
                            preview_cols = st.columns(min(3, len(images)))
                            for i, col in enumerate(preview_cols):
                                if i < len(images):
                                    with col:
                                        st.image(images[i], caption=f"Page {i+1}", use_container_width=True)
                        
                        # Prepare the prompt for the AI
                        input_content = description
                        if images:
                            input_content += f"\n\nThe document contains {len(images)} pages."
                        
                        # Check if we got meaningful content
                        is_extraction_failed = extraction_method.lower() in ["failed", "error"]
                        
                        # Additional context for the model
                        context_message = f"""
                        This is a document that has been processed using {extraction_method}.
                        {'' if not is_extraction_failed else 'WARNING: The extraction process had issues. Create notes based on what information is available.'}
                        Please analyze the content thoroughly and extract ALL important information,
                        focusing on key terms, important sentences, and main concepts.
                        Ensure your notes are comprehensive but concise, and highlight the most critical information.
                        """
                        
                        # Generate notes using the Gemini API correctly
                        model = genai.GenerativeModel(model_name)
                        
                        # If we have images, we can use them directly withGemini (for PDFs)
                        if images and len(images) > 0:
                            # For multi-turn conversation with images
                            chat = model.start_chat(history=[])
                            
                            # Process up to first 5 images (to avoid token limits)
                            image_inputs = []
                            for i, img in enumerate(images[:5]):
                                # Convert PIL image to bytes
                                img_byte_arr = io.BytesIO()
                                img.save(img_byte_arr, format='PNG')
                                img_bytes = img_byte_arr.getvalue()
                                
                                # Add to inputs
                                image_inputs.append({
                                    "mime_type": "image/png",
                                    "data": img_bytes
                                })
                            
                            # Send images with context to the model
                            response = chat.send_message(
                                content=[
                                    selected_prompt,
                                    context_message,
                                    *image_inputs
                                ]
                            )
                            
                            notes_content = response.text
                        else:
                            # Text-only approach for PPTX or failed PDF extraction
                            prompt_with_context = f"{selected_prompt}\n\n{context_message}\n\nDocument content:\n{input_content}"
                            
                            # Check if extraction failed or returned empty/error content
                            if is_extraction_failed or not input_content or input_content.strip() == "":
                                prompt_with_context += "\n\nNOTE: The document extraction failed. Please acknowledge this in your notes and explain what information is missing."
                            
                            response = model.generate_content(prompt_with_context)
                            notes_content = response.text
                        
                        # Store the generated content in session state
                        st.session_state.notes_content = notes_content
                        st.session_state.note_type = note_type
                        
                        # Initialize chat history with the notes as context
                        st.session_state.chat_history = [
                            {"role": "user", "parts": [f"Here are the notes from the document: {notes_content}"]},
                            {"role": "model", "parts": ["Understood, I can answer questions about these notes."]}
                        ]
                        
                    except Exception as e:
                        st.error(f"Error generating notes: {str(e)}")
                        if show_debug:
                            import traceback
                            st.error(traceback.format_exc())
        else:
            st.warning("Please upload a document first.")

    # Display generated notes if they exist in session state
    if st.session_state.notes_content:
        st.markdown(f'<div class="sub-header">{st.session_state.note_type} Generated</div>', unsafe_allow_html=True)
        
        st.markdown(st.session_state.notes_content)
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Export button that always appears when notes are available
        st.markdown(get_download_link(st.session_state.notes_content, st.session_state.note_type), unsafe_allow_html=True)

# Chat Tab - new functionality for chatting about the generated notes
with tab2:
    if st.session_state.notes_content:
        st.markdown('<div class="sub-header">Chat About Your Document Notes</div>', unsafe_allow_html=True)
        with st.form(key='chat_form'):
            question = st.text_input("Ask a question about the notes:")
            submit_button = st.form_submit_button("Ask")
        
        if submit_button and question:
            try:
                model = genai.GenerativeModel(model_name)
                chat = model.start_chat(history=st.session_state.chat_history)
                response = chat.send_message(question)
                answer = response.text
                st.session_state.chat_history.append({"role": "user", "parts": [question]})
                st.session_state.chat_history.append({"role": "model", "parts": [answer]})
            except Exception as e:
                st.error(f"Error: {str(e)}")
        
        # Display conversation history, skipping initial context
        for msg in st.session_state.chat_history[2:]:
            if msg["role"] == "user":
                st.markdown(f"**Q: {msg['parts'][0]}**")
            else:
                st.markdown(f"A: {msg['parts'][0]}")
    else:
        st.info("Please generate notes in the PDF Notes tab first.")

# Adding footer with helpful information
st.markdown("---")
st.markdown("""
### How to get the best results:
- Upload clear, readable PDFs or PPTX files
- For academic papers, technical documents, or reports, try the "Official Notes"
- For general content, use "English Notes" for simplicity
- For multi-language audiences, try "Hinglish Notes"
- Install pdf2image and Poppler for better PDF processing
- If you encounter errors with pdf2image, try installing PyPDF2 as an alternative: `pip install PyPDF2`

#### Installation Requirements:
1. Install Python packages:
   ```
   pip install streamlit pdf2image PyPDF2 google-generativeai python-dotenv pillow python-pptx
   ```
2. Install Poppler:
   - Windows: Download from [poppler-windows](https://github.com/oschwartz10612/poppler-windows/releases/)
   - macOS: `brew install poppler`
   - Linux: `apt-get install poppler-utils`

#### Troubleshooting:
- If you see "Document stream is empty" errors, your document might be corrupted, password-protected, or in an unsupported format
- Try converting your document to a standard format using online converters or Adobe Acrobat before uploading
- Check if your document is password-protected and remove the password before uploading
""")